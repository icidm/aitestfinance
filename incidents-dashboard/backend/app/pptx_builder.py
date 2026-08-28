from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


def _inner_build(incidents, services, lang="en"):
    """Generate the incident presentation using the corporate deck style
    (16:9, Calibri, blue/lilac palette)."""
    txt = {
        "en": {
            "cover_eyebrow": "INCIDENT REPORT - PRODUCTION",
            "cover_title": "Incident Dashboard",
            "cover_desc": "Production incident status: severity, resolution, and impacted services.",
            "cover_generated": "GENERATED",
            "chip_open": "OPEN",
            "chip_services": "SERVICES",
            "chip_resolved": "RESOLVED",
            "summary_eyebrow": "EXECUTIVE SUMMARY",
            "summary_title": "Key indicators",
            "summary_subtitle": "Overview of incident status for the selected period.",
            "kpi_total": "Total incidents",
            "kpi_resolved": "{count} resolved",
            "kpi_open": "Open incidents",
            "kpi_in_progress": "{count} in progress",
            "kpi_mttr": "Mean time to resolution",
            "kpi_mttr_sub": "Average MTTR",
            "kpi_critical": "Open critical",
            "kpi_critical_sub": "{count} critical total",
            "summary_callout": "Resolution rate is {rate}% - {open_count} incidents need active follow-up.",
            "dist_eyebrow": "DISTRIBUTION",
            "dist_title": "Incidents by severity",
            "dist_subtitle": "Breakdown of total incidents by criticality level.",
            "sev_critical": "Critical",
            "sev_high": "High",
            "sev_medium": "Medium",
            "sev_low": "Low",
            "dist_callout": "Total analyzed: {total} incidents - {critical} critical.",
            "timeline_eyebrow": "TIMELINE",
            "timeline_title": "Incidents over time",
            "timeline_subtitle": "Incidents created in the last {days} days.",
            "timeline_no_data": "No incidents in the last {days} days.",
            "table_date": "DATE",
            "table_total": "TOTAL",
            "table_by_severity": "BY SEVERITY",
            "table_open": "OPEN",
            "table_resolved": "RESOLVED",
            "services_eyebrow": "SERVICES",
            "services_title": "Most affected services",
            "services_subtitle": "Top 10 services by incident volume.",
            "table_service": "SERVICE",
            "table_incidents": "INCIDENTS",
            "table_pct_total": "% OF TOTAL",
            "table_weight": "RELATIVE WEIGHT",
            "detail_eyebrow": "DETAIL",
            "detail_title": "Recent incidents",
            "detail_subtitle": "The latest {count} incidents.",
            "table_id": "ID",
            "table_title": "TITLE",
            "table_service_short": "SERVICE",
            "table_sev": "SEV.",
            "table_status": "STATUS",
            "status_open": "Open",
            "status_in_progress": "In progress",
            "status_resolved": "Resolved",
            "sev_abbr_high": "H",
            "sev_abbr_low": "L",
        },
        "es": {
            "cover_eyebrow": "INFORME DE INCIDENCIAS - PRODUCCION",
            "cover_title": "Panel de Incidencias",
            "cover_desc": "Estado de incidencias en produccion: severidad, resolucion y servicios afectados.",
            "cover_generated": "GENERADO",
            "chip_open": "ABIERTAS",
            "chip_services": "SERVICIOS",
            "chip_resolved": "RESUELTAS",
            "summary_eyebrow": "RESUMEN EJECUTIVO",
            "summary_title": "Indicadores clave",
            "summary_subtitle": "Vision general del estado de incidencias del periodo seleccionado.",
            "kpi_total": "Total de incidencias",
            "kpi_resolved": "{count} resueltas",
            "kpi_open": "Casos abiertos",
            "kpi_in_progress": "{count} en progreso",
            "kpi_mttr": "Tiempo medio de resolucion",
            "kpi_mttr_sub": "MTTR promedio",
            "kpi_critical": "Criticas abiertas",
            "kpi_critical_sub": "{count} criticas totales",
            "summary_callout": "La tasa de resolucion es {rate}% - {open_count} incidencias requieren seguimiento activo.",
            "dist_eyebrow": "DISTRIBUCION",
            "dist_title": "Incidencias por severidad",
            "dist_subtitle": "Reparto del total de incidencias segun su criticidad.",
            "sev_critical": "Critica",
            "sev_high": "Alta",
            "sev_medium": "Media",
            "sev_low": "Baja",
            "dist_callout": "Total analizado: {total} incidencias - {critical} criticas.",
            "timeline_eyebrow": "EVOLUCION",
            "timeline_title": "Linea temporal",
            "timeline_subtitle": "Incidencias creadas en los ultimos {days} dias.",
            "timeline_no_data": "Sin incidencias en los ultimos {days} dias.",
            "table_date": "FECHA",
            "table_total": "TOTAL",
            "table_by_severity": "POR SEVERIDAD",
            "table_open": "ABIERTAS",
            "table_resolved": "RESUELTAS",
            "services_eyebrow": "SERVICIOS",
            "services_title": "Servicios mas afectados",
            "services_subtitle": "Top 10 de servicios por volumen de incidencias.",
            "table_service": "SERVICIO",
            "table_incidents": "INCIDENCIAS",
            "table_pct_total": "% DEL TOTAL",
            "table_weight": "PESO RELATIVO",
            "detail_eyebrow": "DETALLE",
            "detail_title": "Incidencias recientes",
            "detail_subtitle": "Las ultimas {count} incidencias.",
            "table_id": "ID",
            "table_title": "TITULO",
            "table_service_short": "SERVICIO",
            "table_sev": "SEV.",
            "table_status": "ESTADO",
            "status_open": "Abierta",
            "status_in_progress": "En progreso",
            "status_resolved": "Resuelta",
            "sev_abbr_high": "A",
            "sev_abbr_low": "B",
        },
    }
    lang = "es" if lang == "es" else "en"

    def tr(key, **kwargs):
        template = txt[lang].get(key, txt["en"].get(key, key))
        return template.format(**kwargs)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ----- Paleta del deck de referencia -----
    ACCENT = RGBColor(0x55, 0x60, 0xE8)  # azul/violeta corporativo
    INK = RGBColor(0x1E, 0x2A, 0x3B)  # títulos
    BODY = RGBColor(0x57, 0x63, 0x7A)  # cuerpo
    MUTED = RGBColor(0x8A, 0x94, 0xA8)  # secundario / eyebrow
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    CARD = RGBColor(0xFF, 0xFF, 0xFF)  # tarjeta blanca
    CARD_ALT = RGBColor(0xFB, 0xFB, 0xFE)  # tarjeta destacada (lila muy claro)
    LILAC = RGBColor(0xEE, 0xF0, 0xFB)  # chip / relleno suave
    LILAC_2 = RGBColor(0xF2, 0xF4, 0xF8)  # chip inactivo
    BAND = RGBColor(0xF4, 0xF5, 0xFC)  # banda callout inferior
    HAIRLINE = RGBColor(0xE4, 0xE7, 0xF0)  # bordes sutiles
    ROW_ALT = RGBColor(0xF7, 0xF8, 0xFC)  # fila alterna tabla

    # Colores de severidad (tonos sobrios integrados con la paleta)
    SEV = {
        "critical": RGBColor(0xE1, 0x3D, 0x5B),
        "high": RGBColor(0xE8, 0x7A, 0x3D),
        "medium": RGBColor(0xE0, 0xB0, 0x3D),
        "low": RGBColor(0x3D, 0xA8, 0x6E),
    }
    SUCCESS = RGBColor(0x2F, 0x9E, 0x64)

    FONT = "Calibri"
    PAGE_W = 13.333

    # ----- Cálculo de métricas -----
    total_incidents = len(incidents)
    open_count = sum(1 for i in incidents if i["status"] in ("open", "in_progress"))
    resolved_count = sum(1 for i in incidents if i["status"] == "resolved")
    in_progress_count = sum(1 for i in incidents if i["status"] == "in_progress")
    critical_open = sum(
        1
        for i in incidents
        if i["severity"] == "critical" and i["status"] in ("open", "in_progress")
    )

    resolved_times = []
    for i in incidents:
        if i["status"] == "resolved" and i.get("resolved_at") and i.get("created_at"):
            diff = (
                datetime.fromisoformat(i["resolved_at"]) - datetime.fromisoformat(i["created_at"])
            ).total_seconds() / 60
            if diff > 0:
                resolved_times.append(diff)
    mttr = round(sum(resolved_times) / len(resolved_times), 1) if resolved_times else 0
    resolution_rate = round(resolved_count / total_incidents * 100) if total_incidents else 0

    by_severity = {"critical": 0, "high": 0, "medium": 0, "low": 0}
    by_status = {"open": 0, "in_progress": 0, "resolved": 0}
    for i in incidents:
        by_severity[i["severity"]] = by_severity.get(i["severity"], 0) + 1
        by_status[i["status"]] = by_status.get(i["status"], 0) + 1

    # ================= Helpers de dibujo =================
    def _no_line(shape):
        shape.line.fill.background()

    def _txt(
        slide,
        left,
        top,
        width,
        height,
        text,
        size,
        color,
        bold=False,
        align=PP_ALIGN.LEFT,
        font=FONT,
        spacing=None,
        anchor=MSO_ANCHOR.TOP,
    ):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
        p.alignment = align
        if spacing is not None:
            # letter-spacing en unidades de 1/100 pt
            run.font._rPr.set("spc", str(int(spacing * 100)))
        return box

    def _spaced(text):
        """Añade tracking manual estilo eyebrow: 'A B C' (espacio simple entre
        letras, triple entre palabras) igual que el deck de referencia."""
        return "   ".join(" ".join(list(w)) for w in text.split(" "))

    def _rect(slide, left, top, width, height, fill=None, line=None, line_w=0.75):
        shp = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        if fill is not None:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        else:
            shp.fill.background()
        if line is not None:
            shp.line.color.rgb = line
            shp.line.width = Pt(line_w)
        else:
            _no_line(shp)
        shp.shadow.inherit = False
        return shp

    def _eyebrow(slide, text, top=0.42, left=0.62):
        _txt(
            slide, left, top, PAGE_W - left - 0.6, 0.3, _spaced(text.upper()), 11, ACCENT, bold=True
        )

    def _header(slide, eyebrow, title, subtitle=""):
        """Encabezado estándar: eyebrow + título grande + subtítulo."""
        _eyebrow(slide, eyebrow, top=0.42)
        _txt(slide, 0.6, 0.72, PAGE_W - 1.2, 0.7, title, 30, INK, bold=True)
        if subtitle:
            _txt(slide, 0.62, 1.36, PAGE_W - 1.2, 0.4, subtitle, 13.5, BODY)

    def _new_slide():
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE
        return s

    def _card(slide, left, top, width, height, fill=CARD, accent_bar=False):
        c = _rect(slide, left, top, width, height, fill=fill, line=HAIRLINE, line_w=0.75)
        if accent_bar:
            _rect(slide, left, top, width, 0.06, fill=ACCENT)
        return c

    def _cell(cell, text, size=10, bold=False, color=INK, align=PP_ALIGN.LEFT, bg=None):
        cell.text = ""
        cell.text_frame.word_wrap = True
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(10)
        cell.margin_right = Pt(10)
        cell.margin_top = Pt(3)
        cell.margin_bottom = Pt(3)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
        p.alignment = align
        if bg is not None:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
        else:
            cell.fill.background()

    date_label = _format_date(datetime.now(), lang)

    # ================= SLIDE 1: PORTADA =================
    s1 = _new_slide()
    _txt(s1, 0.85, 1.95, 11.0, 0.35, _spaced(tr("cover_eyebrow")), 12, ACCENT, bold=True)
    _txt(s1, 0.8, 2.35, 11.5, 1.75, tr("cover_title"), 44, INK, bold=True)
    _txt(s1, 0.85, 3.75, 10.5, 0.8, tr("cover_desc"), 15, BODY)
    _txt(
        s1,
        0.85,
        5.45,
        9.0,
        0.3,
        _spaced(f"{tr('cover_generated')} · {date_label.upper()}"),
        10.5,
        MUTED,
    )
    # Chips resumen a la derecha
    chip_defs = [
        (str(total_incidents), "TOTAL"),
        (str(open_count), tr("chip_open")),
        (str(len(services)), tr("chip_services")),
        (f"{resolution_rate}%", tr("chip_resolved")),
    ]
    cx, cy, cw, ch, gap = 10.35, 2.55, 1.12, 1.12, 0.14
    for k, (val, lab) in enumerate(chip_defs):
        col = k % 2
        row = k // 2
        x = cx + col * (cw + gap)
        y = cy + row * (ch + gap)
        _rect(s1, x, y, cw, ch, fill=LILAC)
        _txt(s1, x, y + 0.18, cw, 0.5, val, 24, ACCENT, bold=True, align=PP_ALIGN.CENTER)
        _txt(s1, x, y + 0.72, cw, 0.25, lab, 8, MUTED, bold=True, align=PP_ALIGN.CENTER)

    # ================= SLIDE 2: RESUMEN EJECUTIVO =================
    s2 = _new_slide()
    _header(s2, tr("summary_eyebrow"), tr("summary_title"), tr("summary_subtitle"))
    kpis = [
        (str(total_incidents), tr("kpi_total"), tr("kpi_resolved", count=resolved_count), True),
        (str(open_count), tr("kpi_open"), tr("kpi_in_progress", count=in_progress_count), False),
        (f"{mttr:.0f} min" if mttr else "—", tr("kpi_mttr"), tr("kpi_mttr_sub"), False),
        (
            str(critical_open),
            tr("kpi_critical"),
            tr("kpi_critical_sub", count=by_severity["critical"]),
            False,
        ),
    ]
    kx, ky, kw, kh, kgap = 0.6, 2.1, 2.95, 2.6, 0.18
    for idx, (val, lab, sub, hot) in enumerate(kpis):
        x = kx + idx * (kw + kgap)
        _card(s2, x, ky, kw, kh, fill=CARD_ALT if hot else CARD, accent_bar=True)
        _txt(s2, x + 0.28, ky + 0.4, kw - 0.5, 0.9, val, 40, INK, bold=True)
        _txt(
            s2,
            x + 0.28,
            ky + 1.42,
            kw - 0.5,
            0.5,
            _spaced(lab.upper()) if len(lab) < 22 else lab.upper(),
            9.5,
            MUTED,
            bold=True,
        )
        _rect(s2, x + 0.28, ky + 1.95, kw - 0.56, 0.02, fill=HAIRLINE)
        _txt(s2, x + 0.28, ky + 2.05, kw - 0.5, 0.4, sub, 10, BODY)
    # Banda callout inferior
    _rect(s2, 0.6, 5.35, PAGE_W - 1.2, 0.72, fill=BAND)
    _txt(
        s2,
        0.9,
        5.35,
        PAGE_W - 1.6,
        0.72,
        tr("summary_callout", rate=resolution_rate, open_count=open_count),
        12,
        BODY,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # ================= SLIDE 3: DISTRIBUCIÓN POR SEVERIDAD =================
    s3 = _new_slide()
    _header(s3, tr("dist_eyebrow"), tr("dist_title"), tr("dist_subtitle"))
    sev_rows = [
        (tr("sev_critical"), by_severity["critical"], SEV["critical"]),
        (tr("sev_high"), by_severity["high"], SEV["high"]),
        (tr("sev_medium"), by_severity["medium"], SEV["medium"]),
        (tr("sev_low"), by_severity["low"], SEV["low"]),
    ]
    max_count = max([c for _, c, _ in sev_rows]) or 1
    bar_x, bar_max_w = 2.3, 8.2
    for idx, (label, count, color) in enumerate(sev_rows):
        y = 2.35 + idx * 0.95
        _txt(
            s3,
            0.6,
            y - 0.02,
            1.5,
            0.4,
            label,
            12,
            INK,
            bold=True,
            align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _rect(s3, bar_x, y, bar_max_w, 0.36, fill=LILAC_2)
        w = (count / max_count) * bar_max_w
        if count > 0:
            _rect(s3, bar_x, y, max(w, 0.08), 0.36, fill=color)
        pct = round(count / total_incidents * 100) if total_incidents else 0
        _txt(
            s3,
            bar_x + bar_max_w + 0.2,
            y - 0.02,
            1.6,
            0.4,
            f"{count}   ({pct}%)",
            13,
            INK,
            bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    _rect(s3, 0.6, 6.25, PAGE_W - 1.2, 0.62, fill=BAND)
    _txt(
        s3,
        0.9,
        6.25,
        PAGE_W - 1.6,
        0.62,
        tr("dist_callout", total=total_incidents, critical=by_severity["critical"]),
        12,
        BODY,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # ================= SLIDE 4: TIMELINE =================
    s4 = _new_slide()
    window_days = 14
    _header(
        s4, tr("timeline_eyebrow"), tr("timeline_title"), tr("timeline_subtitle", days=window_days)
    )
    buckets = {}
    parsed_dates = []
    for i in incidents:
        try:
            parsed_dates.append(datetime.fromisoformat(i["created_at"]))
        except (ValueError, KeyError):
            continue

    if parsed_dates:
        latest_day = max(parsed_dates).replace(hour=0, minute=0, second=0, microsecond=0)
        cutoff_day = latest_day - timedelta(days=window_days - 1)
        end_exclusive = latest_day + timedelta(days=1)

        d = cutoff_day
        while d <= latest_day:
            day_key = d.strftime("%Y-%m-%d")
            buckets[day_key] = {
                "date": day_key,
                "total": 0,
                "critical": 0,
                "high": 0,
                "medium": 0,
                "low": 0,
                "resolved": 0,
            }
            d += timedelta(days=1)

        for i in incidents:
            try:
                created = datetime.fromisoformat(i["created_at"])
            except (ValueError, KeyError):
                continue
            if created < cutoff_day or created >= end_exclusive:
                continue
            day = created.strftime("%Y-%m-%d")
            b = buckets.get(day)
            if not b:
                continue
            b["total"] += 1
            b[i["severity"]] = b.get(i["severity"], 0) + 1
            if i["status"] == "resolved":
                b["resolved"] += 1

    timeline_rows = sorted(buckets.values(), key=lambda x: x["date"])[-10:]
    if timeline_rows:
        headers = [
            tr("table_date"),
            tr("table_total"),
            tr("table_by_severity"),
            tr("table_open"),
            tr("table_resolved"),
        ]
        widths = [1.9, 1.4, 4.5, 2.3, 2.03]
        table = _add_styled_table(
            s4, len(timeline_rows) + 1, headers, widths, 2.05, ACCENT, WHITE, INK, FONT
        )
        for r_idx, row in enumerate(timeline_rows, 1):
            bg = ROW_ALT if r_idx % 2 == 0 else WHITE
            d = datetime.strptime(row["date"], "%Y-%m-%d")
            _cell(table.cell(r_idx, 0), _format_day(d, lang), size=10, bold=True, color=INK, bg=bg)
            _cell(
                table.cell(r_idx, 1),
                row["total"],
                size=11,
                bold=True,
                color=ACCENT,
                bg=bg,
                align=PP_ALIGN.CENTER,
            )
            sev_text = f"C {row['critical']}   {tr('sev_abbr_high')} {row['high']}   M {row['medium']}   {tr('sev_abbr_low')} {row['low']}"
            _cell(table.cell(r_idx, 2), sev_text, size=9.5, color=BODY, bg=bg)
            _cell(
                table.cell(r_idx, 3),
                row["total"] - row["resolved"],
                size=10,
                color=BODY,
                bg=bg,
                align=PP_ALIGN.CENTER,
            )
            _cell(
                table.cell(r_idx, 4),
                row["resolved"],
                size=11,
                bold=True,
                color=SUCCESS,
                bg=bg,
                align=PP_ALIGN.CENTER,
            )
    else:
        _txt(s4, 0.6, 3.0, PAGE_W - 1.2, 0.5, tr("timeline_no_data", days=window_days), 13, MUTED)

    # ================= SLIDE 5: SERVICES =================
    s5 = _new_slide()
    _header(s5, tr("services_eyebrow"), tr("services_title"), tr("services_subtitle"))
    svc_counts = {}
    for i in incidents:
        svc_counts[i["service"]] = svc_counts.get(i["service"], 0) + 1
    svc_sorted = sorted(svc_counts.items(), key=lambda x: x[1], reverse=True)[:10]
    if svc_sorted:
        max_svc = svc_sorted[0][1] or 1
        headers = [
            tr("table_service"),
            tr("table_incidents"),
            tr("table_pct_total"),
            tr("table_weight"),
        ]
        widths = [4.2, 1.9, 1.9, 4.13]
        table = _add_styled_table(
            s5, len(svc_sorted) + 1, headers, widths, 2.05, ACCENT, WHITE, INK, FONT
        )
        for r_idx, (sname, scount) in enumerate(svc_sorted, 1):
            bg = ROW_ALT if r_idx % 2 == 0 else WHITE
            pct = (scount / total_incidents * 100) if total_incidents else 0
            _cell(table.cell(r_idx, 0), sname, size=10.5, bold=True, color=INK, bg=bg)
            _cell(
                table.cell(r_idx, 1),
                scount,
                size=11,
                bold=True,
                color=ACCENT,
                bg=bg,
                align=PP_ALIGN.CENTER,
            )
            _cell(
                table.cell(r_idx, 2),
                f"{pct:.1f}%",
                size=10,
                color=BODY,
                bg=bg,
                align=PP_ALIGN.CENTER,
            )
            # mini-barra dentro de la celda de peso: se dibuja como shape encima
            _cell(table.cell(r_idx, 3), "", bg=bg)
        # Barras de peso relativo (dibujadas sobre la última columna)
        col_x = 0.6 + sum(widths[:3])
        col_w = widths[3]
        row_h = 4.6 / len(svc_sorted) if len(svc_sorted) else 0.4
        top_tbl = 2.05
        header_h = 0.45
        for r_idx, (sname, scount) in enumerate(svc_sorted):
            frac = scount / max_svc
            y = top_tbl + header_h + r_idx * row_h + row_h / 2 - 0.09
            _rect(s5, col_x + 0.15, y, col_w - 0.45, 0.18, fill=LILAC_2)
            _rect(s5, col_x + 0.15, y, max((col_w - 0.45) * frac, 0.06), 0.18, fill=ACCENT)

    # ================= SLIDE 6: RECENT INCIDENTS =================
    s6 = _new_slide()
    shown = min(18, total_incidents)
    _header(s6, tr("detail_eyebrow"), tr("detail_title"), tr("detail_subtitle", count=shown))
    display_incidents = sorted(incidents, key=lambda x: x.get("created_at", ""), reverse=True)[:18]
    if display_incidents:
        headers = [
            tr("table_id"),
            tr("table_title"),
            tr("table_service_short"),
            tr("table_sev"),
            tr("table_status"),
        ]
        widths = [0.8, 5.6, 2.6, 1.0, 2.13]
        table = _add_styled_table(
            s6, len(display_incidents) + 1, headers, widths, 2.05, ACCENT, WHITE, INK, FONT
        )
        status_map = {
            "open": tr("status_open"),
            "in_progress": tr("status_in_progress"),
            "resolved": tr("status_resolved"),
        }
        sev_map = {
            "critical": "C",
            "high": tr("sev_abbr_high"),
            "medium": "M",
            "low": tr("sev_abbr_low"),
        }
        for r_idx, inc in enumerate(display_incidents, 1):
            bg = ROW_ALT if r_idx % 2 == 0 else WHITE
            sev_color = SEV.get(inc["severity"], MUTED)
            st_color = SUCCESS if inc["status"] == "resolved" else BODY
            _cell(
                table.cell(r_idx, 0),
                f"#{inc['id']}",
                size=9,
                color=MUTED,
                bg=bg,
                align=PP_ALIGN.CENTER,
            )
            _cell(table.cell(r_idx, 1), inc["title"][:64], size=9.5, color=INK, bg=bg)
            _cell(table.cell(r_idx, 2), inc["service"][:28], size=9, color=BODY, bg=bg)
            _cell(
                table.cell(r_idx, 3),
                sev_map.get(inc["severity"], "?"),
                size=10,
                bold=True,
                color=sev_color,
                bg=bg,
                align=PP_ALIGN.CENTER,
            )
            _cell(
                table.cell(r_idx, 4),
                status_map.get(inc["status"], inc["status"]),
                size=9,
                bold=True,
                color=st_color,
                bg=bg,
                align=PP_ALIGN.CENTER,
            )

    # ----- Footers con numeración total -----
    total_pages = len(prs.slides)
    lang_badge = "EN" if lang == "en" else "ES"
    for idx, slide in enumerate(prs.slides, 1):
        if idx == 1:
            continue  # la portada no lleva footer
        _rect(slide, PAGE_W - 2.55, 7.03, 0.5, 0.22, fill=LILAC, line=HAIRLINE, line_w=0.6)
        _txt(
            slide,
            PAGE_W - 2.55,
            7.03,
            0.5,
            0.22,
            lang_badge,
            8.5,
            ACCENT,
            bold=True,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _txt(
            slide,
            PAGE_W - 1.55,
            7.08,
            0.95,
            0.3,
            f"{idx:02d} / {total_pages:02d}",
            10,
            MUTED,
            align=PP_ALIGN.RIGHT,
        )

    return prs


def _add_styled_table(slide, n_rows, headers, widths, top, accent, header_fg, ink, font):
    """Crea una tabla con cabecera de acento y devuelve el objeto table."""
    left = 0.6
    total_w = sum(widths)
    table_shape = slide.shapes.add_table(
        n_rows, len(headers), Inches(left), Inches(top), Inches(total_w), Inches(4.6)
    )
    table = table_shape.table
    # Desactiva el banding por defecto de PowerPoint
    tblPr = table._tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        tblPr.set("firstRow", "0")
        tblPr.set("bandRow", "0")
    for c_idx, w in enumerate(widths):
        table.columns[c_idx].width = Inches(w)
    table.rows[0].height = Inches(0.45)
    for c_idx, htext in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.text = ""
        cell.text_frame.word_wrap = True
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(10)
        cell.margin_right = Pt(10)
        cell.margin_top = Pt(3)
        cell.margin_bottom = Pt(3)
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = htext
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = header_fg
        run.font.name = font
        p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
    return table


_MONTHS_EN = [
    "January",
    "February",
    "March",
    "April",
    "May",
    "June",
    "July",
    "August",
    "September",
    "October",
    "November",
    "December",
]
_MONTHS_EN_ABBR = [
    "Jan",
    "Feb",
    "Mar",
    "Apr",
    "May",
    "Jun",
    "Jul",
    "Aug",
    "Sep",
    "Oct",
    "Nov",
    "Dec",
]
_DAYS_EN_ABBR = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]
_MONTHS_ES = [
    "enero",
    "febrero",
    "marzo",
    "abril",
    "mayo",
    "junio",
    "julio",
    "agosto",
    "septiembre",
    "octubre",
    "noviembre",
    "diciembre",
]
_MONTHS_ES_ABBR = [
    "ene",
    "feb",
    "mar",
    "abr",
    "may",
    "jun",
    "jul",
    "ago",
    "sep",
    "oct",
    "nov",
    "dic",
]
_DAYS_ES_ABBR = ["lun", "mar", "mie", "jue", "vie", "sab", "dom"]


def _format_date(dt, lang="en"):
    if lang == "es":
        return f"{dt.day} de {_MONTHS_ES[dt.month - 1]} de {dt.year}"
    return f"{dt.day} {_MONTHS_EN[dt.month - 1]} {dt.year}"


def _format_day(dt, lang="en"):
    if lang == "es":
        return f"{_DAYS_ES_ABBR[dt.weekday()]} {dt.day} {_MONTHS_ES_ABBR[dt.month - 1]}"
    return f"{_DAYS_EN_ABBR[dt.weekday()]} {dt.day} {_MONTHS_EN_ABBR[dt.month - 1]}"


def build_pptx(incidents, services, lang="en"):
    return _inner_build(incidents, services, lang)
